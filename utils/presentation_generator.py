import os
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage, PageBreak, Frame, PageTemplate
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from reportlab.pdfgen import canvas
from datetime import datetime
import re
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)

class PresentationGenerator:
    """Generate eye-catching technical presentations - 1 page per item"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.setup_custom_styles()
    
    def setup_custom_styles(self):
        """Setup custom paragraph styles for presentations"""
        self.title_style = ParagraphStyle(
            'PresentationTitle',
            parent=self.styles['Heading1'],
            fontSize=28,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=20,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        self.item_title_style = ParagraphStyle(
            'ItemTitle',
            parent=self.styles['Heading1'],
            fontSize=22,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=15,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        self.spec_heading_style = ParagraphStyle(
            'SpecHeading',
            fontSize=14,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=8,
            fontName='Helvetica-Bold'
        )
        
        self.spec_text_style = ParagraphStyle(
            'SpecText',
            fontSize=11,
            textColor=colors.black,
            spaceAfter=6,
            alignment=TA_JUSTIFY
        )

    def _get_logo_path(self):
        candidates = [
            os.path.join('static', 'images', 'al-shaya-logo-white@2x.png'),  # White logo first for presentations
            os.path.join('static', 'images', 'AlShaya-Logo-color@2x.png'),
            os.path.join('static', 'images', 'LOGO.png')
        ]
        for p in candidates:
            if os.path.exists(p):
                return p
        return None

    def _get_white_logo_path(self):
        """Get white logo specifically for PPTX presentations"""
        white_logo = os.path.join('static', 'images', 'al-shaya-logo-white@2x.png')
        if os.path.exists(white_logo):
            return white_logo
        return self._get_logo_path()  # Fallback to default

    def _draw_header_footer(self, canv: canvas.Canvas, doc):
        """Draw properly placed header logo and footer website for presentation PDF."""
        page_width, page_height = doc.pagesize
        gold = colors.HexColor('#d4af37')
        dark = colors.HexColor('#1a365d')
        
        # Header gold line
        canv.setStrokeColor(gold)
        canv.setLineWidth(2)
        canv.line(doc.leftMargin, page_height - 40, page_width - doc.rightMargin, page_height - 40)
        
        # Logo centered in header - larger and more visible
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                w, h = 130, 46  # Larger logo
                x = (page_width - w) / 2  # Center horizontally
                y = page_height - 38
                canv.drawImage(logo, x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
            except Exception:
                pass
        
        # Footer with gold line and website centered
        canv.setStrokeColor(gold)
        canv.setLineWidth(2)
        canv.line(doc.leftMargin, doc.bottomMargin + 15, page_width - doc.rightMargin, doc.bottomMargin + 15)
        
        canv.setFillColor(dark)
        canv.setFont('Helvetica', 10)
        footer_text = 'https://alshayaenterprises.com'
        canv.drawCentredString(page_width / 2, doc.bottomMargin + 5, footer_text)
    
    def generate(self, file_id, session, format_type='pdf'):
        """
        Generate technical presentation with 1 page/slide per item
        Always generates PPTX first, then converts to PDF if needed
        Args:
            file_id: The file ID
            session: Flask session
            format_type: 'pdf' or 'pptx'
        Returns: path to generated file (PDF or PPTX), and stores PPTX path in session
        """
        # Get file info and extracted data
        uploaded_files = session.get('uploaded_files', [])
        file_info = None
        
        for f in uploaded_files:
            if f['id'] == file_id:
                file_info = f
                break
        
        if not file_info:
            raise Exception('File not found. Please upload and extract a file first.')
        
        # Get costed data (preferred) or stitched table or extraction result
        if 'costed_data' in file_info:
            items = self.parse_items_from_costed_data(file_info['costed_data'], session, file_id)
        elif 'stitched_table' in file_info:
            items = self.parse_items_from_stitched_table(file_info['stitched_table'], session, file_id)
        elif 'extraction_result' in file_info:
            items = self.parse_items_from_extraction(file_info['extraction_result'], session, file_id)
        else:
            raise Exception('No data available. Please extract tables first.')
        
        if not items:
            raise Exception('No items found in the table.')
        
        # Create output directory
        session_id = session['session_id']
        output_dir = os.path.join('outputs', session_id, 'presentations')
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate file based on format
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # ALWAYS generate PPTX first (will be available for download)
        pptx_file = os.path.join(output_dir, f'presentation_{file_id}_{timestamp}.pptx')
        self.generate_pptx(items, pptx_file)
        
        # Store PPTX path in file_info for later download
        file_info['presentation_pptx'] = pptx_file
        
        if format_type == 'pptx':
            output_file = pptx_file
        else:  # pdf - convert PPTX to PDF to maintain same layout
            output_file = os.path.join(output_dir, f'presentation_{file_id}_{timestamp}.pdf')
            self.convert_pptx_to_pdf(pptx_file, output_file)
        
        return output_file
    
    def parse_items_from_costed_data(self, costed_data, session, file_id):
        """Parse items from costed table data"""
        items = []
        session_id = session.get('session_id', '')
        
        for table in costed_data.get('tables', []):
            headers = [h for h in table.get('headers', []) if str(h).lower() not in ['action', 'actions', 'product selection', 'productselection']]
            
            for row in table.get('rows', []):
                # Find description column
                description = ''
                for h in headers:
                    h_str = str(h).lower() if h else ''
                    if 'descript' in h_str or 'discript' in h_str or 'item' in h_str:  # Handle misspelling
                        description = self.strip_html(row.get(h, ''))
                        break
                
                # Find quantity
                qty = ''
                unit = ''
                for h in headers:
                    h_str = str(h).lower() if h else ''
                    if 'qty' in h_str or 'quantity' in h_str:
                        qty = self.strip_html(row.get(h, ''))
                    if 'unit' in h_str and 'rate' not in h_str:
                        unit = self.strip_html(row.get(h, ''))
                
                # Find pricing
                unit_rate = ''
                total = ''
                for h in headers:
                    h_str = str(h).lower() if h else ''
                    if 'rate' in h_str or 'price' in h_str:
                        unit_rate = self.strip_html(row.get(h, ''))
                    if 'total' in h_str or 'amount' in h_str:
                        total = self.strip_html(row.get(h, ''))
                
                # Find image
                image_path = None
                for h in headers:
                    cell_value = row.get(h, '')
                    if self.contains_image(cell_value):
                        image_path = self.extract_image_path(cell_value, session_id, file_id)
                        break
                
                if description:  # Only add if we have a description
                    item = {
                        'description': description,
                        'qty': qty,
                        'unit': unit,
                        'unit_rate': unit_rate,
                        'total': total,
                        'image_path': image_path,
                        'brand': self.extract_brand(description),
                        'specifications': self.extract_specifications(description)
                    }
                    items.append(item)
        
        return items
    
    def parse_items_from_stitched_table(self, stitched_table, session, file_id):
        """Parse items from stitched HTML table data"""
        items = []
        session_id = session.get('session_id', '')
        
        # Parse the HTML
        html_content = stitched_table.get('html', '')
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Find the table
        table = soup.find('table')
        if not table:
            logger.error("No table found in stitched HTML")
            return items
        
        # Get headers
        headers = []
        header_row = table.find('tr')
        if header_row:
            for th in header_row.find_all(['th', 'td']):
                header_text = th.get_text(strip=True).lower()
                # Exclude Product Selection and Actions columns
                if header_text not in ['action', 'actions', 'product selection', 'productselection']:
                    headers.append(header_text)
        
        logger.info(f"Found headers: {headers}")
        logger.info(f"Checking for description header in: {headers}")
        
        # Get data rows (skip header row)
        rows = table.find_all('tr')[1:]  # Skip first row (headers)
        logger.info(f"Found {len(rows)} data rows")
        
        for row in rows:
            cells = row.find_all('td')
            
            # Build row dict, skipping Product Selection and Actions cells
            row_data = {}
            col_idx = 0
            for i, cell in enumerate(cells):
                # Skip Product Selection and Actions cells
                if cell.find(class_='product-selection-dropdowns') or cell.find('button'):
                    continue
                text = cell.get_text(strip=True).lower()
                if 'product selection' in text or 'actions' in text:
                    continue
                
                if col_idx < len(headers):
                    # Check if cell contains image
                    img = cell.find('img')
                    if img:
                        row_data[headers[col_idx]] = str(cell)  # Keep HTML with image
                    else:
                        row_data[headers[col_idx]] = cell.get_text(strip=True)
                    col_idx += 1
            
            # Extract fields
            description = ''
            logger.debug(f"Looking for description in headers: {headers}")
            logger.debug(f"Row data keys: {list(row_data.keys())}")
            
            for h in headers:
                # Ensure header is a string
                h_str = str(h).lower() if h else ''
                if 'descript' in h_str or 'item' in h_str or 'discript' in h_str:  # Handle misspelling
                    description = self.strip_html(row_data.get(h, ''))
                    logger.debug(f"Found description in header '{h}': {description[:50] if description else 'empty'}")
                    break
            
            if not description:
                logger.warning(f"No description found in row. Headers: {headers}, Row data keys: {list(row_data.keys())}")
                continue
            
            # Find quantity
            qty = ''
            unit = ''
            for h in headers:
                h_str = str(h).lower() if h else ''
                if 'qty' in h_str or 'quantity' in h_str:
                    qty = self.strip_html(row_data.get(h, ''))
                if 'unit' in h_str and 'rate' not in h_str and 'price' not in h_str:
                    unit = self.strip_html(row_data.get(h, ''))
            
            # Find pricing
            unit_rate = ''
            total = ''
            for h in headers:
                h_str = str(h).lower() if h else ''
                if ('rate' in h_str or 'price' in h_str) and 'unit' in h_str:
                    unit_rate = self.strip_html(row_data.get(h, ''))
                if 'total' in h_str or 'amount' in h_str:
                    total = self.strip_html(row_data.get(h, ''))
            
            # Find image
            image_path = None
            for h in headers:
                cell_value = row_data.get(h, '')
                if self.contains_image(str(cell_value)):
                    image_path = self.extract_image_path(str(cell_value), session_id, file_id)
                    break
            
            item = {
                'description': description,
                'qty': qty,
                'unit': unit,
                'unit_rate': unit_rate,
                'total': total,
                'image_path': image_path,
                'brand': self.extract_brand(description),
                'specifications': self.extract_specifications(description)
            }
            items.append(item)
        
        logger.info(f"Parsed {len(items)} items from stitched table")
        return items
    
    def strip_html(self, text):
        """Strip HTML tags from text"""
        return re.sub(r'<[^>]+>', '', str(text)).strip()
    
    def contains_image(self, cell_value):
        """Check if cell contains an image reference"""
        return '<img' in str(cell_value).lower() or 'img_in_' in str(cell_value).lower()
    
    def extract_image_path(self, cell_value, session_id, file_id):
        """Extract image path from cell value"""
        try:
            match = re.search(r'src=["\']([^"\']+)["\']', str(cell_value))
            if match:
                img_path = match.group(1).lstrip('/')
                
                # Handle URLs (http/https)
                if img_path.startswith('http://') or img_path.startswith('https://'):
                    return img_path
                
                # Handle local paths - if already starts with outputs, return as-is
                if img_path.startswith('outputs'):
                    # Check if file exists
                    if os.path.exists(img_path):
                        return img_path
                    # If not, might be missing file - still return it
                    return img_path
                else:
                    # Ensure all parts are strings
                    if isinstance(session_id, (list, tuple)):
                        session_id = session_id[0] if session_id else ''
                    if isinstance(file_id, (list, tuple)):
                        file_id = file_id[0] if file_id else ''
                    if isinstance(img_path, (list, tuple)):
                        img_path = img_path[0] if img_path else ''
                    
                    # Check if it's a relative path that needs to be joined
                    full_path = os.path.join('outputs', str(session_id), str(file_id), str(img_path))
                    if os.path.exists(full_path):
                        return full_path
                    # Also try without the session_id/file_id prefix in case it's already included
                    if os.path.exists(str(img_path)):
                        return str(img_path)
                    return full_path  # Return even if doesn't exist yet, let download logic handle it
            
            # Look for any image path pattern in imgs/ folder (more flexible regex)
            if 'imgs/' in str(cell_value):
                match = re.search(r'(imgs/[^"\s<>]+\.(jpg|png|jpeg|gif|webp))', str(cell_value), re.IGNORECASE)
                if match:
                    img_relative_path = match.group(1)
                    # Ensure all parts are strings
                    if isinstance(session_id, (list, tuple)):
                        session_id = session_id[0] if session_id else ''
                    if isinstance(file_id, (list, tuple)):
                        file_id = file_id[0] if file_id else ''
                    full_path = os.path.join('outputs', str(session_id), str(file_id), str(img_relative_path))
                    if os.path.exists(full_path):
                        return full_path
                    return full_path  # Return even if doesn't exist yet
        except Exception as e:
            logger.error(f"Error extracting image path: {e}")
            pass
        return None
    
    def generate_pdf(self, items, output_file):
        """Generate PDF presentation"""
        doc = SimpleDocTemplate(output_file, pagesize=A4, 
                                topMargin=1.0*inch, bottomMargin=0.8*inch,
                                leftMargin=0.75*inch, rightMargin=0.75*inch)
        story = []
        
        # Cover page
        story.extend(self.create_cover_page())
        story.append(PageBreak())
        
        # Create one page per item
        for idx, item in enumerate(items):
            story.extend(self.create_item_page_pdf(item, idx + 1))
            if idx < len(items) - 1:
                story.append(PageBreak())
        
        # Build PDF with header/footer
        doc.build(story, onFirstPage=self._draw_header_footer, onLaterPages=self._draw_header_footer)
    
    def generate_pptx(self, items, output_file):
        """Generate PowerPoint presentation"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        self.create_title_slide_pptx(prs)
        
        # Add one slide per item
        for idx, item in enumerate(items):
            self.create_item_slide_pptx(prs, item, idx + 1)
        
        prs.save(output_file)
    
    def convert_pptx_to_pdf(self, pptx_file, pdf_file):
        """Convert PPTX to PDF using PowerPoint COM automation on Windows"""
        try:
            import platform
            import time
            
            if platform.system() == 'Windows':
                try:
                    import comtypes.client
                    
                    logger.info(f"Converting PPTX to PDF using PowerPoint COM: {pptx_file}")
                    
                    # Convert paths to absolute
                    pptx_abs = os.path.abspath(pptx_file)
                    pdf_abs = os.path.abspath(pdf_file)
                    
                    logger.info(f"Absolute paths - PPTX: {pptx_abs}, PDF: {pdf_abs}")
                    
                    # Create PowerPoint application
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    powerpoint.Visible = 1  # Must be visible for COM automation
                    
                    logger.info("PowerPoint application created")
                    
                    # Open presentation without window restrictions
                    presentation = powerpoint.Presentations.Open(pptx_abs)
                    
                    logger.info("Presentation opened")
                    
                    # Save as PDF (32 = ppSaveAsPDF)
                    presentation.SaveAs(pdf_abs, 32)
                    
                    logger.info("SaveAs called")
                    
                    # Close and cleanup
                    presentation.Close()
                    powerpoint.Quit()
                    
                    logger.info("PowerPoint closed")
                    
                    # Wait for file to be fully written
                    max_wait = 10  # 10 seconds max
                    wait_time = 0
                    while wait_time < max_wait:
                        if os.path.exists(pdf_abs) and os.path.getsize(pdf_abs) > 1000:
                            logger.info(f"PDF created successfully: {pdf_abs} ({os.path.getsize(pdf_abs)} bytes)")
                            return
                        time.sleep(0.5)
                        wait_time += 0.5
                    
                    logger.warning(f"PDF file not found or too small after {max_wait} seconds")
                    
                except Exception as e:
                    logger.error(f"PowerPoint COM conversion failed: {e}")
                    import traceback
                    logger.error(traceback.format_exc())
                    raise
            else:
                raise Exception("PowerPoint COM automation is only available on Windows")
            
        except Exception as e:
            logger.error(f"Error converting PPTX to PDF: {e}")
            raise Exception(f"Could not convert presentation to PDF. Please ensure Microsoft PowerPoint is installed. Error: {str(e)}")
    
    def create_title_slide_pptx(self, prs):
        """Create PowerPoint title slide with enhanced design"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background - Navy blue header box (increased height to contain full logo)
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), 
            Inches(10), Inches(3.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(26, 54, 93)  # Navy blue
        header_shape.line.fill.background()
        
        # Gold accent bar (moved down to be below the logo)
        accent_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(3.2),
            Inches(10), Inches(0.15)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = RGBColor(212, 175, 55)  # Gold
        accent_shape.line.fill.background()
        
        # Logo centered in header with more vertical space - use white logo
        logo = self._get_white_logo_path()
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, Inches(3.5), Inches(0.8), width=Inches(3))
            except Exception:
                pass
        
        # Title with navy background (no company name - removed as requested)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "TECHNICAL PROPOSAL"
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(48)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(26, 54, 93)  # Navy
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(5.0), Inches(6), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Furniture, Fixtures & Equipment"
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(26)
        subtitle_p.font.color.rgb = RGBColor(100, 116, 139)  # Gray
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(3), Inches(6.2), Inches(4), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime('%B %d, %Y')
        date_p = date_frame.paragraphs[0]
        date_p.font.size = Pt(18)
        date_p.font.color.rgb = RGBColor(71, 85, 105)  # Dark gray
        date_p.alignment = PP_ALIGN.CENTER
        
        # Footer with website
        footer_box = slide.shapes.add_textbox(Inches(2.5), Inches(6.8), Inches(5), Inches(0.4))
        footer_frame = footer_box.text_frame
        footer_frame.text = "https://alshayaenterprises.com"
        footer_p = footer_frame.paragraphs[0]
        footer_p.font.size = Pt(14)
        footer_p.font.color.rgb = RGBColor(212, 175, 55)  # Gold
        footer_p.alignment = PP_ALIGN.CENTER
    
    def create_item_slide_pptx(self, prs, item, page_num):
        """Create PowerPoint slide for one item with enhanced design"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Header bar with navy blue background - increased height
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(1.1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(26, 54, 93)  # Navy blue
        header_shape.line.fill.background()
        
        # Gold accent line under header
        accent_line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(1.1),
            Inches(10), Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = RGBColor(212, 175, 55)  # Gold
        accent_line.line.fill.background()
        
        # Small white logo top-right in header
        logo = self._get_white_logo_path()
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, Inches(8.2), Inches(0.2), width=Inches(1.5))
            except Exception:
                pass
        
        # Title in header - positioned to not overlap logo (left-aligned with more space)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7.2), Inches(0.6))
        title_frame = title_box.text_frame
        title_text = f"Item {page_num}: {item['description'][:55]}"
        title_frame.text = title_text
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(20)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text on navy
        
        # Image (left side) - adjusted position to account for taller header
        image_path = item.get('image_path')
        if image_path:
            # If it's a URL, download it first
            if image_path.startswith('http'):
                from utils.image_helper import download_image
                cached_path = download_image(image_path)
                if cached_path:
                    image_path = cached_path
            
            if image_path and os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(image_path, Inches(0.6), Inches(1.8), 
                                            width=Inches(4.2), height=Inches(4.2))
                except Exception as e:
                    # If image fails, add placeholder
                    img_placeholder = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(4.2), Inches(1))
                    img_frame = img_placeholder.text_frame
                    img_frame.text = "[Image Not Available]"
                    img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                # Add placeholder
                img_placeholder = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(4.2), Inches(1))
                img_frame = img_placeholder.text_frame
                img_frame.text = "[Image Not Available]"
                img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        else:
            # Add placeholder
            img_placeholder = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(4.2), Inches(1))
            img_frame = img_placeholder.text_frame
            img_frame.text = "[Image Not Available]"
            img_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Details box (right side) - adjusted position for taller header
        details_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.0))
        details_frame = details_box.text_frame
        details_frame.word_wrap = True
        
        # Product Details heading
        p = details_frame.paragraphs[0]
        p.text = "Product Details"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 54, 93)  # Navy
        p.space_after = Pt(12)
        
        # Full description paragraph
        desc_p = details_frame.add_paragraph()
        desc_p.text = item['description']
        desc_p.font.size = Pt(13)
        desc_p.font.color.rgb = RGBColor(51, 51, 51)  # Dark text
        desc_p.space_after = Pt(10)
        
        # Key Details - professional text only (no icons)
        p = details_frame.add_paragraph()
        p.text = f"Brand: {item['brand']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.space_after = Pt(6)
        
        p = details_frame.add_paragraph()
        p.text = f"Quantity: {item['qty']} {item['unit']}"
        p.font.size = Pt(13)
        p.space_after = Pt(12)
        
        # Specifications heading
        p = details_frame.add_paragraph()
        p.text = "Specifications:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 54, 93)
        p.space_after = Pt(8)
        
        # Add specifications (limit to prevent overflow)
        for spec in item['specifications'][:6]:  # Limit to 6 specs
            p = details_frame.add_paragraph()
            p.text = f"• {spec}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.level = 1
            p.space_after = Pt(3)
        
        # Footer with website
        footer_box = slide.shapes.add_textbox(Inches(3), Inches(7), Inches(4), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = "https://alshayaenterprises.com"
        footer_p = footer_frame.paragraphs[0]
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(100, 116, 139)
        footer_p.alignment = PP_ALIGN.CENTER
    
    def create_item_page_pdf(self, item, page_num):
        """Create PDF page for one item"""
        story = []
        
        # Item title
        item_title = f"Item {page_num}: {item['description'][:80]}"
        story.append(Paragraph(item_title, self.item_title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # Create two-column layout
        left_content = []
        right_content = []
        
        # Left: Image
        image_path = item.get('image_path')
        if image_path:
            # If it's a URL, download it first
            if image_path.startswith('http'):
                from utils.image_helper import download_image
                cached_path = download_image(image_path)
                if cached_path:
                    image_path = cached_path
            
            if image_path and os.path.exists(image_path):
                try:
                    img = RLImage(image_path, width=2.5*inch, height=2.5*inch)
                    left_content.append(img)
                except Exception as e:
                    left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
            else:
                left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        else:
            left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        
        # Right: Details
        details_html = f"""
            <para>
                <b><font size="14" color="#1a365d">Product Details</font></b><br/>
                <br/>
                <b>Brand:</b> {item['brand']}<br/>
                <b>Quantity:</b> {item['qty']} {item['unit']}<br/>
                <b>Unit Rate:</b> {item['unit_rate']}<br/>
                <b>Total Amount:</b> {item['total']}<br/>
                <br/>
                <b><font color="#1a365d">Specifications:</font></b><br/>
            </para>
        """
        right_content.append(Paragraph(details_html, self.spec_text_style))
        
        # Specifications
        for spec in item['specifications']:
            right_content.append(Paragraph(f"• {spec}", self.spec_text_style))
            right_content.append(Spacer(1, 0.05*inch))
        
        # Two-column table
        data = [[left_content, right_content]]
        t = Table(data, colWidths=[3*inch, 3.5*inch])
        t.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 10),
            ('RIGHTPADDING', (0, 0), (-1, -1), 10),
        ]))
        
        story.append(t)
        
        return story
    
    def parse_items_from_extraction(self, extraction_result, session, file_id):
        """
        Parse individual items from extraction result
        Returns: list of item dictionaries
        """
        items = []
        
        for layout_result in extraction_result.get('layoutParsingResults', []):
            markdown_text = layout_result.get('markdown', {}).get('text', '')
            images = layout_result.get('markdown', {}).get('images', {})
            
            # Parse tables from markdown
            table_rows = self.extract_table_rows(markdown_text)
            
            for row in table_rows:
                item = {
                    'sn': row.get('sn', row.get('sl.no', '')),
                    'description': row.get('description', row.get('item', '')),
                    'qty': row.get('qty', row.get('quantity', '')),
                    'unit': row.get('unit', ''),
                    'unit_rate': row.get('unit rate', row.get('unit price', row.get('rate', ''))),
                    'total': row.get('total', row.get('amount', '')),
                    'image': self.find_item_image(row, images),
                    'brand': self.extract_brand(row.get('description', '')),
                    'specifications': self.extract_specifications(row.get('description', ''))
                }
                items.append(item)
        
        return items
    
    def extract_table_rows(self, markdown_text):
        """Extract table rows from markdown text"""
        lines = markdown_text.split('\n')
        rows = []
        headers = []
        
        for line in lines:
            if '|' in line:
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                
                if not headers and cells:
                    # First row is headers
                    headers = [h.lower() for h in cells]
                elif headers and cells and not all(c in '-: ' for c in ''.join(cells)):
                    # Data row
                    if len(cells) == len(headers):
                        row = dict(zip(headers, cells))
                        rows.append(row)
        
        return rows
    
    def find_item_image(self, row, images):
        """Find image associated with this item"""
        # Try to find image reference in row
        for key, value in row.items():
            if 'image' in str(key).lower():
                # Check if value references an image
                value_str = str(value) if value else ''
                for img_path, img_url in images.items():
                    img_path_str = str(img_path) if img_path else ''
                    img_url_str = str(img_url) if img_url else ''
                    if value_str in img_path_str or img_path_str in value_str:
                        return img_url_str
                
        # Return first available image if no specific match
        if images:
            first_img = list(images.values())[0]
            return str(first_img) if first_img else None
        
        return None
    
    def extract_brand(self, description):
        """Extract brand name from description (simple heuristic)"""
        # Common brand patterns - this is simplified
        brands = ['Sedus', 'Narbutas', 'Sokoa', 'B&T', 'Herman Miller', 'Steelcase', 'Haworth', 'Knoll']
        
        for brand in brands:
            if brand.lower() in description.lower():
                return brand
        
        # Try to extract first capitalized word
        words = description.split()
        for word in words:
            if word and word[0].isupper() and len(word) > 2:
                return word
        
        return 'Premium Brand'
    
    def extract_specifications(self, description):
        """Extract specifications from description"""
        # Split description into bullet points
        specs = []
        
        # Look for dimensions
        dimension_pattern = r'\d+\s*[xX×]\s*\d+\s*[xX×]?\s*\d*\s*(mm|cm|m|inch|in|")'
        dimensions = re.findall(dimension_pattern, description)
        if dimensions:
            specs.append(f"Dimensions: {', '.join(dimensions)}")
        
        # Look for materials
        materials = ['wood', 'metal', 'steel', 'aluminum', 'fabric', 'leather', 'plastic', 'glass', 'laminate']
        found_materials = [mat for mat in materials if mat in description.lower()]
        if found_materials:
            specs.append(f"Materials: {', '.join(found_materials).title()}")
        
        # Look for colors
        colors_list = ['black', 'white', 'grey', 'gray', 'brown', 'blue', 'red', 'green', 'beige']
        found_colors = [col for col in colors_list if col in description.lower()]
        if found_colors:
            specs.append(f"Available Colors: {', '.join(found_colors).title()}")
        
        if not specs:
            # Use description as-is if no specific specs found
            specs.append(description[:200])
        
        return specs
    
    def create_cover_page(self):
        """Create presentation cover page"""
        story = []
        
        # Centered large logo (if available)
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                img = RLImage(logo, width=3.5*inch, height=3.5*inch)
                img.hAlign = 'CENTER'
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
            except Exception:
                pass

        # Title
        title = Paragraph("TECHNICAL PROPOSAL", self.title_style)
        story.append(Spacer(1, 2*inch))
        story.append(title)
        story.append(Spacer(1, 0.5*inch))
        
        # Subtitle
        subtitle_style = ParagraphStyle(
            'Subtitle',
            parent=self.styles['Normal'],
            fontSize=16,
            textColor=colors.HexColor('#1a365d'),
            alignment=TA_CENTER
        )
        subtitle = Paragraph("Furniture, Fixtures & Equipment", subtitle_style)
        story.append(subtitle)
        story.append(Spacer(1, 1*inch))
        
        # Company info
        company_info = f"""
            <para align="center">
                <b>Prepared By:</b><br/>
                <font size="14" color="#667eea"><b>Your Company Name</b></font><br/>
                <br/>
                Date: {datetime.now().strftime('%B %d, %Y')}<br/>
            </para>
        """
        story.append(Paragraph(company_info, self.styles['Normal']))
        
        return story
    
    def create_item_page(self, item, page_num):
        """Create one page for an item with eye-catching design"""
        story = []
        
        # Item number and title
        item_title = f"Item {page_num}: {item['description'][:60]}"
        story.append(Paragraph(item_title, self.item_title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Create two-column layout using table
        left_content = []
        right_content = []
        
        # Left column - Image
        if item['image']:
            try:
                # For now, placeholder - in production, download and embed image
                img_placeholder = Paragraph(
                    f'<para align="center"><b>[Product Image]</b><br/>{item["image"][:50]}...</para>',
                    self.styles['Normal']
                )
                left_content.append(img_placeholder)
            except:
                left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        else:
            left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        
        # Right column - Specifications
        specs_html = f"""
            <para>
                <b><font size="14" color="#667eea">Product Details</font></b><br/>
                <br/>
                <b>Brand:</b> {item['brand']}<br/>
                <b>Quantity:</b> {item['qty']} {item['unit']}<br/>
                <b>Unit Rate:</b> {item['unit_rate']}<br/>
                <b>Total Amount:</b> {item['total']}<br/>
                <br/>
                <b><font color="#667eea">Specifications:</font></b><br/>
            </para>
        """
        right_content.append(Paragraph(specs_html, self.spec_text_style))
        
        # Add specifications as bullet points
        for spec in item['specifications']:
            spec_bullet = f"• {spec}"
            right_content.append(Paragraph(spec_bullet, self.spec_text_style))
            right_content.append(Spacer(1, 0.1*inch))
        
        # Additional info
        additional_info = """
            <para>
                <br/>
                <b><font color="#667eea">Additional Information:</font></b><br/>
                • Country of Origin: Various<br/>
                • Warranty: As per manufacturer standard<br/>
                • Lead Time: 4-6 weeks<br/>
                • Finish: As specified or equivalent<br/>
            </para>
        """
        right_content.append(Paragraph(additional_info, self.spec_text_style))
        
        # Create two-column table
        data = [[left_content, right_content]]
        
        col_widths = [3*inch, 3.5*inch]
        t = Table(data, colWidths=col_widths)
        t.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 10),
            ('RIGHTPADDING', (0, 0), (-1, -1), 10),
        ]))
        
        story.append(t)
        story.append(Spacer(1, 0.3*inch))
        
        # Bottom section - Key features
        features_title = Paragraph('<b><font size="12" color="#1a365d">KEY FEATURES</font></b>', self.styles['Normal'])
        story.append(features_title)
        story.append(Spacer(1, 0.1*inch))
        
        features = [
            "✓ Premium quality construction",
            "✓ Modern ergonomic design",
            "✓ Environmentally friendly materials",
            "✓ Easy maintenance and durability"
        ]
        
        features_text = '<br/>'.join(features)
        story.append(Paragraph(features_text, self.spec_text_style))
        
        return story
